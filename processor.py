from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re
import fitz  # PyMuPDF

# Variable globale pour stocker le texte extrait
extracted_text = ""

def creer_presentation_ppt(data, output_file):
    # Créer une nouvelle présentation
    prs = Presentation()

    # Définir les styles de texte
    title_style = prs.slide_layouts[0]
    content_style = prs.slide_layouts[1]

    # Ajouter une diapositive de titre
    slide = prs.slides.add_slide(title_style)
    title = slide.shapes.title
    title.text = "Algeria Tech Infographie"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Ajouter une diapositive de contenu pour chaque section
    for section, items in data.items():
        slide = prs.slides.add_slide(content_style)
        title = slide.shapes.title
        title.text = section
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Créer une table pour les données
        rows = len(items) + 1
        cols = 3
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(0.8 * rows)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Définir les en-têtes de colonne
        table.cell(0, 0).text = "Indicateur"
        table.cell(0, 1).text = "Valeur"
        table.cell(0, 2).text = "Unité"

        # Remplir la table avec les données
        for i, (key, value) in enumerate(items.items(), start=1):
            # Nettoyer les données
            cleaned_value = nettoyer_donnees(value)

            table.cell(i, 0).text = key
            table.cell(i, 1).text = str(cleaned_value)
            table.cell(i, 2).text = obtenir_unite(key)

            # Appliquer le style de texte
            for j in range(cols):
                cell = table.cell(i, j)
                cell.text_frame.paragraphs[0].font.size = Pt(24)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Appliquer le style de fond de tableau
        for row in table.rows:
            for cell in row.cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(25, 25, 112)  # Bleu foncé

    # Sauvegarder la présentation
    prs.save(output_file)

def nettoyer_donnees(value):
    # Filtrer les années
    if isinstance(value, (int, float)) and 2000 <= value <= 2030:
        return ""

    # Supprimer les .0 et convertir en int
    if isinstance(value, float) and value.is_integer():
        return int(value)

    return value

def obtenir_unite(key):
    # Associer les bonnes unités
    unite_map = {
        "abonnés": "Abonnés",
        "téléchargements": "Go",
        "stockage": "To",
        "budget": "DA"
    }

    for k, v in unite_map.items():
        if k in key.lower():
            return v

    return ""

def analyser_pdf():
    global extracted_text
    pdf_path = input("Veuillez glisser-déposer le PDF ici et appuyez sur Entrée : ")
    pdf_path = pdf_path.strip('"')  # Nettoyer le chemin

    try:
        doc = fitz.open(pdf_path)
        extracted_text = ""
        for page in doc:
            extracted_text += page.get_text()
        print("✅ PDF analysé avec succès !")
    except Exception as e:
        print(f"❌ Erreur lors de l'analyse du PDF : {e}")

def generer_data_show():
    global extracted_text
    if not extracted_text:
        print("❌ Aucun texte extrait. Veuillez d'abord analyser un PDF.")
        return

    print("⏳ Génération du Data Show en cours...")
    # Ici, vous pouvez ajouter le code pour traiter le texte extrait et générer le Data Show
    # Pour l'instant, nous utilisons les données d'exemple
    data = {
        "Télécommunications": {
            "Abonnés mobiles": 45000000,
            "Téléchargements mobiles": 1200000000000,
            "Budget": 1500000000
        },
        "Technologie": {
            "Entreprises tech": 2500,
            "Startups": 1500,
            "Budget R&D": 500000000
        }
    }
    creer_presentation_ppt(data, "presentation.pptx")
    print("✅ Data Show généré avec succès !")

def main():
    while True:
        print("\nMenu Principal:")
        print("1. Analyser PDF (Synthèse)")
        print("2. Générer Data Show (Depuis le PDF)")
        print("3. Quitter")

        choix = input("Entrez votre choix (1-3): ")

        if choix == "1":
            analyser_pdf()
        elif choix == "2":
            generer_data_show()
        elif choix == "3":
            print("Au revoir !")
            break
        else:
            print("Choix invalide. Veuillez entrer un nombre entre 1 et 3.")

if __name__ == "__main__":
    main()